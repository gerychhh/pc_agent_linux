from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from tools import filesystem


def _result(ok: bool, **kwargs: Any) -> str:
    payload = {"ok": ok, **kwargs}
    if not ok and "error" not in payload:
        payload["error"] = "unknown_error"
    return json.dumps(payload, ensure_ascii=False)


def _verify_file(path: str) -> dict[str, Any]:
    target = Path(path)
    verified = target.exists() and target.is_file() and target.stat().st_size > 0
    return {
        "verified": verified,
        "verify_reason": "file_exists_and_nonempty" if verified else "file_missing_or_empty",
    }


def _load_docx() -> Any:
    from docx import Document  # type: ignore

    return Document


def _load_openpyxl() -> Any:
    import openpyxl  # type: ignore

    return openpyxl


def _load_pptx() -> Any:
    from pptx import Presentation  # type: ignore

    return Presentation


def create_docx(path: str, title: str | None, paragraphs: list[str]) -> str:
    raw = filesystem.create_docx(path, title, paragraphs)
    try:
        parsed = json.loads(raw)
    except json.JSONDecodeError:
        parsed = {"ok": False, "error": "invalid_result", "path": path}
    verification = _verify_file(path)
    return _result(
        bool(parsed.get("ok")),
        path=parsed.get("path", path),
        error=parsed.get("error"),
        verified=verification["verified"],
        verify_reason=verification["verify_reason"],
    )


def read_docx(path: str, max_paragraphs: int = 200) -> str:
    try:
        document = _load_docx()(path)
        paragraphs = [p.text for p in document.paragraphs if p.text]
        limited = paragraphs[: int(max_paragraphs)]
        verification = _verify_file(path)
        return _result(
            True,
            path=path,
            text="\n".join(limited),
            paragraphs=len(limited),
            verified=verification["verified"],
            verify_reason=verification["verify_reason"],
            error=None,
        )
    except Exception as exc:
        verification = _verify_file(path)
        return _result(
            False,
            path=path,
            text="",
            paragraphs=0,
            verified=verification["verified"],
            verify_reason=verification["verify_reason"],
            error=str(exc),
        )


def append_docx(path: str, paragraphs: list[str]) -> str:
    try:
        Document = _load_docx()
        target = Path(path)
        if target.exists():
            document = Document(path)
        else:
            document = Document()
        for paragraph in paragraphs:
            document.add_paragraph(paragraph)
        target.parent.mkdir(parents=True, exist_ok=True)
        document.save(str(target))
        verification = _verify_file(path)
        return _result(
            True,
            path=str(target),
            added=len(paragraphs),
            verified=verification["verified"],
            verify_reason=verification["verify_reason"],
            error=None,
        )
    except Exception as exc:
        verification = _verify_file(path)
        return _result(
            False,
            path=path,
            added=0,
            verified=verification["verified"],
            verify_reason=verification["verify_reason"],
            error=str(exc),
        )


def replace_in_docx(path: str, find_text: str, replace_text: str) -> str:
    try:
        Document = _load_docx()
        document = Document(path)
        replacements = 0
        for paragraph in document.paragraphs:
            if find_text in paragraph.text:
                paragraph.text = paragraph.text.replace(find_text, replace_text)
                replacements += 1
        document.save(path)
        verification = _verify_file(path)
        return _result(
            True,
            path=path,
            replacements=replacements,
            verified=verification["verified"],
            verify_reason=verification["verify_reason"],
            error=None,
        )
    except Exception as exc:
        verification = _verify_file(path)
        return _result(
            False,
            path=path,
            replacements=0,
            verified=verification["verified"],
            verify_reason=verification["verify_reason"],
            error=str(exc),
        )


def create_xlsx(path: str, sheets: list[dict[str, Any]]) -> str:
    try:
        openpyxl = _load_openpyxl()
        workbook = openpyxl.Workbook()
        if workbook.sheetnames:
            default_sheet = workbook.active
            workbook.remove(default_sheet)
        for sheet in sheets:
            name = sheet.get("name") or "Sheet"
            rows = sheet.get("rows") or []
            worksheet = workbook.create_sheet(title=str(name))
            for row in rows:
                worksheet.append(list(row))
        target = Path(path)
        target.parent.mkdir(parents=True, exist_ok=True)
        workbook.save(str(target))
        verification = _verify_file(path)
        return _result(
            True,
            path=str(target),
            sheets=len(sheets),
            verified=verification["verified"],
            verify_reason=verification["verify_reason"],
            error=None,
        )
    except Exception as exc:
        verification = _verify_file(path)
        return _result(
            False,
            path=path,
            sheets=0,
            verified=verification["verified"],
            verify_reason=verification["verify_reason"],
            error=str(exc),
        )


def read_xlsx(path: str, sheet: str | None = None, max_rows: int = 200) -> str:
    try:
        openpyxl = _load_openpyxl()
        workbook = openpyxl.load_workbook(path, data_only=True)
        worksheet = workbook[sheet] if sheet else workbook.active
        rows = []
        for idx, row in enumerate(worksheet.iter_rows(values_only=True)):
            if idx >= int(max_rows):
                break
            rows.append(list(row))
        verification = _verify_file(path)
        return _result(
            True,
            path=path,
            sheet=worksheet.title,
            rows=rows,
            verified=verification["verified"],
            verify_reason=verification["verify_reason"],
            error=None,
        )
    except Exception as exc:
        verification = _verify_file(path)
        return _result(
            False,
            path=path,
            sheet=sheet,
            rows=[],
            verified=verification["verified"],
            verify_reason=verification["verify_reason"],
            error=str(exc),
        )


def write_xlsx_cell(path: str, sheet: str, cell: str, value: Any) -> str:
    try:
        openpyxl = _load_openpyxl()
        target = Path(path)
        if target.exists():
            workbook = openpyxl.load_workbook(str(target))
        else:
            workbook = openpyxl.Workbook()
        worksheet = workbook[sheet] if sheet in workbook.sheetnames else workbook.create_sheet(title=sheet)
        worksheet[cell] = value
        target.parent.mkdir(parents=True, exist_ok=True)
        workbook.save(str(target))
        verification = _verify_file(path)
        return _result(
            True,
            path=str(target),
            sheet=sheet,
            cell=cell,
            verified=verification["verified"],
            verify_reason=verification["verify_reason"],
            error=None,
        )
    except Exception as exc:
        verification = _verify_file(path)
        return _result(
            False,
            path=path,
            sheet=sheet,
            cell=cell,
            verified=verification["verified"],
            verify_reason=verification["verify_reason"],
            error=str(exc),
        )


def create_pptx(path: str, title: str | None, slides: list[dict[str, Any]]) -> str:
    try:
        Presentation = _load_pptx()
        presentation = Presentation()
        if title:
            layout = presentation.slide_layouts[0]
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = ""
        for slide_data in slides:
            layout = presentation.slide_layouts[1]
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = str(slide_data.get("title") or "")
            body = slide.shapes.placeholders[1].text_frame
            body.clear()
            for bullet in slide_data.get("bullets") or []:
                body.add_paragraph().text = str(bullet)
        target = Path(path)
        target.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(target))
        verification = _verify_file(path)
        return _result(
            True,
            path=str(target),
            slides=len(slides),
            verified=verification["verified"],
            verify_reason=verification["verify_reason"],
            error=None,
        )
    except Exception as exc:
        verification = _verify_file(path)
        return _result(
            False,
            path=path,
            slides=0,
            verified=verification["verified"],
            verify_reason=verification["verify_reason"],
            error=str(exc),
        )


def append_pptx_slide(path: str, title: str | None, bullets: list[str]) -> str:
    try:
        Presentation = _load_pptx()
        presentation = Presentation(path)
        layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title or ""
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for bullet in bullets:
            body.add_paragraph().text = str(bullet)
        target = Path(path)
        target.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(target))
        verification = _verify_file(path)
        return _result(
            True,
            path=str(target),
            verified=verification["verified"],
            verify_reason=verification["verify_reason"],
            error=None,
        )
    except Exception as exc:
        verification = _verify_file(path)
        return _result(
            False,
            path=path,
            verified=verification["verified"],
            verify_reason=verification["verify_reason"],
            error=str(exc),
        )
